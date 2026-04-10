from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
PINK = RGBColor(0xE8, 0x56, 0x6A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xFF, 0xF5, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE5, 0xE5)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=Pt(12), bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline(slide, lines, left, top, width, height, font_size=Pt(11), color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

# ──────────────────────────────────────────────
# SLIDE 1 – TITLE
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# Pink left bar
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)

# White background
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)

# Pink accent block (top right)
add_rect(slide, 9.5, 0, 3.83, 2.5, fill_color=LIGHT)

# Logo / company
add_text(slide, "SUPERSTREUSEL", 1.2, 0.7, 8, 0.6, font_size=Pt(11), bold=True, color=PINK)
add_text(slide, "UG (haftungsbeschränkt)", 1.2, 1.2, 8, 0.4, font_size=Pt(9), color=GRAY)

# Title
add_text(slide, "Wege in die", 1.2, 2.2, 9, 0.8, font_size=Pt(40), bold=False, color=DARK)
add_text(slide, "Selbstständigkeit", 1.2, 2.9, 9, 0.9, font_size=Pt(40), bold=True, color=DARK)

# Subtitle
add_text(slide, "Gründungsschritte · Chancen · Risiken", 1.2, 3.9, 9, 0.5, font_size=Pt(14), color=GRAY)

# Pink line
add_rect(slide, 1.2, 4.6, 4, 0.04, fill_color=PINK)

# Class info
add_text(slide, "Schulprojekt — Klasse · 2026", 1.2, 4.8, 6, 0.4, font_size=Pt(10), color=GRAY)

# Right side note
add_text(slide, "Unternehmen:", 9.7, 0.3, 3.3, 0.3, font_size=Pt(8), color=GRAY)
add_text(slide, "superstreusel.de", 9.7, 0.6, 3.3, 0.3, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Gegründet von\nDani & Kati Decker", 9.7, 1.0, 3.3, 0.5, font_size=Pt(9), color=GRAY)

# Bottom bar
add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Gruppe 6 Personen · Rechtsform: UG · Eigenkapital: 10.000 €", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 2 – UNTERNEHMENSVORSTELLUNG
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "UNTERNEHMENSVORSTELLUNG", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Superstreusel UG — Wer wir sind", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

# Geschäftsidee (left, full width top)
add_rect(slide, 1.0, 1.6, 11.5, 1.2, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
add_text(slide, "Geschäftsidee", 1.2, 1.7, 4, 0.35, font_size=Pt(10), bold=True, color=PINK)
add_text(slide, "Online-Shop für Spezialzuckerstreusel & Backdekorationen — thematisch sortiert nach Anlass (Hochzeit, Geburtstag, Weihnachten, Ostern). Qualität, die im Einzelhandel nicht erhältlich ist.", 1.2, 2.05, 11.0, 0.55, font_size=Pt(11), color=GRAY)

# Three boxes bottom
cols = [
    ("Zielgruppe", "Hobbybäckerinnen (25–45 J.) · Profibäcker & Konditoreien · Eventcaterer\nÜberwiegend weiblich, Instagram-aktiv, Anlass-orientiert"),
    ("Alleinstellungsmerkmal (USP)", "Einziger spezialisierter Online-Shop für thematische Spezialstreusel in Deutschland. Kein Supermarkt-Sortiment — nur handverlesene Qualität. Starke Community & First-Mover-Vorteil."),
    ("Vision & Ziele", "Marktführer für Backdekorationen im D-A-CH-Raum bis 2027.\n500+ Produkte · 100.000 Community-Mitglieder · B2B-Netzwerk aus Konditoreien aufbauen"),
]
for i, (title, body) in enumerate(cols):
    x = 1.0 + i * 3.9
    add_rect(slide, x, 3.0, 3.6, 3.5, fill_color=PINK if i == 1 else LIGHT, line_color=BORDER if i != 1 else None, line_width=Pt(1))
    tc = WHITE if i == 1 else DARK
    gc = WHITE if i == 1 else GRAY
    add_text(slide, title, x + 0.2, 3.15, 3.2, 0.4, font_size=Pt(11), bold=True, color=tc)
    add_rect(slide, x + 0.2, 3.6, 1.0, 0.04, fill_color=WHITE if i == 1 else PINK)
    add_text(slide, body, x + 0.2, 3.75, 3.2, 2.5, font_size=Pt(10), color=gc)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Superstreusel UG (haftungsbeschränkt) · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 3 – AUSGANGSSITUATION
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "AUSGANGSSITUATION", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Die Geschichte hinter Superstreusel", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

# Two column layout
# Left column
add_text(slide, "Der Anfang", 1.0, 1.7, 5.5, 0.4, font_size=Pt(13), bold=True, color=DARK)
add_rect(slide, 1.0, 2.1, 1.5, 0.05, fill_color=PINK)
add_multiline(slide, [
    "Dani Decker backt Kuchen auf Stöckchen — erst für Freunde, dann für Kunden.",
    "",
    "Kein Geschäft, keine Küche, kein Plan. Der Teig kommt aus Krefeld, die Lieferung geht durch Hamburg."
], 1.0, 2.2, 5.5, 2.0, font_size=Pt(11), color=GRAY)

add_text(slide, "Schwester Kati kommt dazu", 1.0, 4.3, 5.5, 0.4, font_size=Pt(13), bold=True, color=DARK)
add_rect(slide, 1.0, 4.7, 1.5, 0.05, fill_color=PINK)
add_multiline(slide, [
    "Beide schreiben Unternehmen an — Weihnachtsfeiern, Jubiläen.",
    "Erste Bestellungen kommen. Das Modell wird professionalisiert."
], 1.0, 4.8, 5.5, 1.5, font_size=Pt(11), color=GRAY)

# Right column
add_text(slide, "Das zweite Geschäftsfeld", 7.2, 1.7, 5.5, 0.4, font_size=Pt(13), bold=True, color=DARK)
add_rect(slide, 7.2, 2.1, 1.5, 0.05, fill_color=PINK)
add_multiline(slide, [
    "Kunden wollen besondere Verzierungen — aber im Supermarkt gibt es nur 0815-Streuware.",
    "",
    "Aus diesem Problem entsteht Superstreusel: ein Online-Shop für Spezialzuckerstreusel."
], 7.2, 2.2, 5.6, 2.0, font_size=Pt(11), color=GRAY)

# Highlight box
add_rect(slide, 7.2, 4.3, 5.6, 1.8, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
add_text(slide, "Startkapital: 10.000 € Eigenkapital", 7.5, 4.5, 5.0, 0.4, font_size=Pt(11), bold=True, color=PINK)
add_text(slide, "Kein einziger Kredit aufgenommen.\nRechtsform: UG (haftungsbeschränkt)", 7.5, 4.9, 5.0, 0.8, font_size=Pt(10), color=GRAY)

# Bottom
add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Quelle: Ausgangssituation – Guter Kuchen e.K. (Schulunterlagen) · superstreusel.de", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 3 – RECHTSFORMWECHSEL
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "PERSON 1 — RECHTSFORMWECHSEL", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Von der Einzelunternehmerin zur UG", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

# Three boxes
boxes = [
    ("Einzelunternehmen", "Nur eine Inhaberin · Keine Haftungsbeschränkung · Kein Mindestkapital", "Start"),
    ("UG (haftungsb.)", "Zwei Gesellschafterinnen möglich · Stammkapital ab 1 € · Haftung beschränkt", "Gewählt"),
    ("GmbH", "Volle Kapitalgesellschaft · 25.000 € Mindestkapital · Höherer Aufwand", "Ziel"),
]

for i, (title, desc, label) in enumerate(boxes):
    x = 1.0 + i * 4.0
    is_chosen = label == "Gewählt"
    bg = PINK if is_chosen else LIGHT
    tc = WHITE if is_chosen else DARK
    gc = WHITE if is_chosen else GRAY
    add_rect(slide, x, 1.7, 3.6, 3.5, fill_color=bg, line_color=BORDER if not is_chosen else None)
    add_text(slide, label.upper(), x + 0.2, 1.9, 3.2, 0.4, font_size=Pt(8), bold=True, color=tc)
    add_text(slide, title, x + 0.2, 2.3, 3.2, 0.5, font_size=Pt(14), bold=True, color=tc)
    add_rect(slide, x + 0.2, 2.85, 1.0, 0.04, fill_color=WHITE if is_chosen else PINK)
    add_multiline(slide, desc.split(" · "), x + 0.2, 3.0, 3.2, 1.8, font_size=Pt(10), color=gc)

add_rect(slide, 1.0, 5.4, 11.5, 0.04, fill_color=BORDER)
add_text(slide, "Warum UG?", 1.0, 5.6, 4, 0.35, font_size=Pt(11), bold=True, color=DARK)
add_multiline(slide, [
    "· Beide Schwestern als gleichberechtigte Gesellschafterinnen",
    "· Haftungsschutz ohne das hohe Mindestkapital der GmbH",
    "· Thesaurierungspflicht: 25% des Gewinns → Rücklage bis 25.000 €"
], 1.0, 5.95, 5.5, 1.0, font_size=Pt(10), color=GRAY)
add_text(slide, "Gründungskosten UG:", 7.0, 5.6, 4, 0.35, font_size=Pt(11), bold=True, color=DARK)
add_multiline(slide, [
    "· Notar: ca. 300–600 €",
    "· Handelsregistereintrag: ca. 150 €",
    "· Laufend: Jahresabschluss + Veröffentlichung"
], 7.0, 5.95, 5.5, 1.0, font_size=Pt(10), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Person 1 · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 5 – GRÜNDUNGSSCHRITTE
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "GRÜNDUNGSSCHRITTE", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Von der Idee zur eingetragenen UG — Schritt für Schritt", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

gsteps = [
    ("1", "Notar & Gesellschaftsvertrag", "Beurkundung des Gesellschaftsvertrags (Name, Sitz Hamburg, Stammkapital 10.000 €, je 50% Anteile). Kosten: ca. 450 €"),
    ("2", "Geschäftskonto eröffnen", "Einzahlung des Stammkapitals auf separates Konto. Bankbescheinigung für Handelsregistereintrag erforderlich."),
    ("3", "Handelsregistereintrag", "Notar leitet Antrag ans Amtsgericht Hamburg weiter. Mit Eintragung entsteht die UG als juristische Person. Ca. 150 €"),
    ("4", "Gewerbeanmeldung", "Anmeldung beim Bezirksamt Hamburg-Mitte. Gewerbeart: Onlinehandel / Einzelhandel Lebensmittel. Ca. 20–60 €"),
    ("5", "Finanzamt & Steuernummer", "Finanzamt erhält Meldung vom Gewerbeamt. Steuerlicher Erfassungsbogen ausfüllen → Steuernummer + USt-IdNr. für EU."),
    ("6", "IHK & Versicherungen", "Pflichtmitgliedschaft IHK Hamburg (automatisch). Betriebshaftpflicht + Produkthaftpflicht abschließen. Ca. 80 €/Monat."),
]

for i, (num, title, desc) in enumerate(gsteps):
    col = i % 3
    row = i // 2
    x = 1.0 + col * 4.0
    y = 1.6 + row * 2.4
    add_rect(slide, x, y, 0.5, 0.5, fill_color=PINK)
    add_text(slide, num, x + 0.12, y + 0.08, 0.3, 0.35, font_size=Pt(14), bold=True, color=WHITE)
    add_rect(slide, x + 0.5, y, 3.1, 0.5, fill_color=LIGHT)
    add_text(slide, title, x + 0.65, y + 0.08, 2.9, 0.35, font_size=Pt(11), bold=True, color=DARK)
    add_text(slide, desc, x + 0.2, y + 0.6, 3.4, 1.6, font_size=Pt(9.5), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Person 3 · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 6 – FINANZIERUNG
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "PERSON 2 — FINANZPLANUNG", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "10.000 € Eigenkapital — kein einziger Kredit", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

# Left: Startkosten Aufstellung
add_text(slide, "Startkosten (einmalig)", 1.0, 1.6, 5.5, 0.35, font_size=Pt(12), bold=True, color=DARK)
add_rect(slide, 1.0, 1.95, 5.5, 0.04, fill_color=PINK)
kosten = [
    ("Notar + Gesellschaftsvertrag", "450 €"),
    ("Handelsregistereintrag", "150 €"),
    ("Küchenausstattung & Geräte", "2.800 €"),
    ("Erster Warenbestand", "1.500 €"),
    ("Verpackungsmaterial", "600 €"),
    ("Website + Online-Shop", "300 €"),
    ("Logo & Verpackungsdesign", "700 €"),
    ("Reserve (Puffer)", "3.500 €"),
]
for i, (item, betrag) in enumerate(kosten):
    y = 2.1 + i * 0.42
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, 1.0, y, 5.5, 0.4, fill_color=bg)
    add_text(slide, item, 1.15, y + 0.07, 3.8, 0.3, font_size=Pt(9.5), color=GRAY)
    add_text(slide, betrag, 5.2, y + 0.07, 1.1, 0.3, font_size=Pt(9.5), bold=True, color=DARK, align=PP_ALIGN.RIGHT)
add_rect(slide, 1.0, 5.5, 5.5, 0.4, fill_color=PINK)
add_text(slide, "Gesamt", 1.15, 5.57, 3.0, 0.3, font_size=Pt(10), bold=True, color=WHITE)
add_text(slide, "10.000 €", 5.2, 5.57, 1.1, 0.3, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Right: Fixkosten + Break-even
add_text(slide, "Monatliche Fixkosten", 7.0, 1.6, 5.8, 0.35, font_size=Pt(12), bold=True, color=DARK)
add_rect(slide, 7.0, 1.95, 5.8, 0.04, fill_color=PINK)
fixkosten = [
    ("Miete Küche/Lager Hamburg", "650 €"),
    ("Shopify + Domain", "30 €"),
    ("Versandmaterial (variabel)", "200 €"),
    ("Steuerberater", "120 €"),
    ("Versicherungen", "80 €"),
]
for i, (item, betrag) in enumerate(fixkosten):
    y = 2.1 + i * 0.42
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, 7.0, y, 5.8, 0.4, fill_color=bg)
    add_text(slide, item, 7.15, y + 0.07, 4.2, 0.3, font_size=Pt(9.5), color=GRAY)
    add_text(slide, betrag, 11.45, y + 0.07, 1.1, 0.3, font_size=Pt(9.5), bold=True, color=DARK, align=PP_ALIGN.RIGHT)
add_rect(slide, 7.0, 4.3, 5.8, 0.4, fill_color=PINK)
add_text(slide, "Fixkosten/Monat", 7.15, 4.37, 3.0, 0.3, font_size=Pt(10), bold=True, color=WHITE)
add_text(slide, "1.080 €", 11.45, 4.37, 1.1, 0.3, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# Break-even box
add_rect(slide, 7.0, 4.9, 5.8, 1.6, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
add_text(slide, "Break-even-Analyse", 7.2, 5.0, 5.4, 0.35, font_size=Pt(11), bold=True, color=PINK)
add_multiline(slide, [
    "Durchschn. Bestellwert: 35 €  ·  Marge: 60%  →  21 € Deckungsbeitrag",
    "Break-even: 1.080 € ÷ 21 € = ca. 52 Bestellungen/Monat",
    "→ Erreicht nach wenigen Wochen durch organische Mundpropaganda",
], 7.2, 5.4, 5.4, 1.0, font_size=Pt(10), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Person 2 · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 5 – PERSONALBEDARF
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "PERSON 3 — PERSONALBEDARF", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Wann einstellen — und wen zuerst?", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

steps = [
    ("Phase 1", "Minijob Versand", "Pakete packen, Ware einlagern.\nErste Entlastung, geringe Fixkosten."),
    ("Phase 2", "Festanstellung\nProduktion", "Wenn tägliches Volumen\neinen vollen Arbeitstag füllt."),
    ("Phase 3", "Freelancer\nMarketing", "Fotografie, Social Media,\nTexte — projektweise."),
    ("Phase 4", "Kundenservice", "Ab ca. 500+ Bestellungen\npro Monat sinnvoll."),
]

for i, (phase, title, desc) in enumerate(steps):
    x = 1.0 + i * 3.0
    add_rect(slide, x, 1.7, 2.7, 0.5, fill_color=PINK)
    add_text(slide, phase, x + 0.15, 1.82, 2.4, 0.3, font_size=Pt(10), bold=True, color=WHITE)
    add_rect(slide, x, 2.2, 2.7, 3.0, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
    add_text(slide, title, x + 0.2, 2.4, 2.4, 0.7, font_size=Pt(12), bold=True, color=DARK)
    add_text(slide, desc, x + 0.2, 3.15, 2.4, 1.4, font_size=Pt(10), color=GRAY)

add_rect(slide, 1.0, 5.5, 11.5, 0.04, fill_color=BORDER)
add_text(slide, "Rechtliche Pflichten bei Festanstellung:", 1.0, 5.7, 6, 0.35, font_size=Pt(11), bold=True, color=DARK)
add_multiline(slide, [
    "· Sozialversicherungsanmeldung (Kranken-, Renten-, Pflegeversicherung)",
    "· Lohnsteueranmeldung · Urlaubsanspruch · Lohnfortzahlung im Krankheitsfall"
], 1.0, 6.1, 11.5, 0.7, font_size=Pt(10), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Person 3 · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 6 – LIEFERANTENSUCHE (Duo A)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "DUO (TEIL A) — LIEFERANTENSUCHE", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Wo findet man Spezialzuckerstreusel-Hersteller?", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

ways = [
    ("Fachmessen", "ISM Köln · Anuga", "Direktkontakt, Muster in der Hand halten, Preise verhandeln. Wichtigster Kanal für Qualitäts-Hersteller."),
    ("B2B-Plattformen", "Alibaba · Global Sources", "Große Auswahl, günstigere Preise. Aber: hohe Mindestmengen, Qualitätskontrolle nötig, immer Muster bestellen."),
    ("Branchennetzwerk", "Innungen · Verbände", "Empfehlungen aus der Praxis. Ein Kollege, der liefert, ist mehr wert als jede Online-Bewertung."),
]

for i, (name, sub, desc) in enumerate(ways):
    x = 1.0 + i * 4.0
    add_rect(slide, x, 1.7, 3.6, 3.8, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
    num = str(i + 1)
    add_rect(slide, x, 1.7, 0.5, 0.5, fill_color=PINK)
    add_text(slide, num, x + 0.12, 1.76, 0.3, 0.35, font_size=Pt(14), bold=True, color=WHITE)
    add_text(slide, name, x + 0.65, 1.78, 2.9, 0.35, font_size=Pt(13), bold=True, color=DARK)
    add_text(slide, sub, x + 0.2, 2.25, 3.2, 0.35, font_size=Pt(10), color=PINK, italic=True)
    add_text(slide, desc, x + 0.2, 2.7, 3.2, 2.5, font_size=Pt(10), color=GRAY)

add_rect(slide, 1.0, 5.7, 11.5, 0.8, fill_color=LIGHT)
add_text(slide, "Wichtig: Immer mindestens 2 Lieferanten pro Kategorie · Vorlaufzeit Sonderstreusel: 3–6 Monate · Langfristige Verträge anstreben", 1.2, 5.85, 11.0, 0.5, font_size=Pt(10), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Duo · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 7 – MARKETINGSTRATEGIE
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "PERSON 4 — MARKETINGSTRATEGIE", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Wie Superstreusel Kunden gewinnt und hält", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

channels = [
    ("Instagram", "Organic + UGC", "Backwelt ist überproportional aktiv. Kunden posten, tagen, teilen — kostenlose Reichweite."),
    ("Pinterest", "Evergreen-Inhalte", "Pins bringen noch nach 2 Jahren Traffic. Backanleitungen + Streusel-Inspiration = idealer Content."),
    ("Google\nShopping", "Saisonal + kaufintentional", "Zu Weihnachten, Ostern, Karneval gezielt schalten. Niedriges Budget, hohe Absicht."),
    ("E-Mail", "Bestandskunden", "Günstigste Reaktivierung. Liste gehört dem Unternehmen — kein Algorithmus kann sie wegnehmen."),
    ("B2B-Direkt", "Konditoreien,\nCaterer", "Größere Bestellungen, regelmäßiger. Musterboxen + Mengenrabatt als Einstieg."),
]

for i, (name, sub, desc) in enumerate(channels):
    col = i % 3
    row = i // 3
    x = 1.0 + col * 4.0
    y = 1.7 + row * 2.5
    add_rect(slide, x, y, 3.6, 2.2, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
    add_text(slide, name, x + 0.2, y + 0.15, 3.2, 0.5, font_size=Pt(13), bold=True, color=DARK)
    add_text(slide, sub, x + 0.2, y + 0.6, 3.2, 0.35, font_size=Pt(9), color=PINK, italic=True)
    add_text(slide, desc, x + 0.2, y + 1.0, 3.2, 1.0, font_size=Pt(9), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Person 4 · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 8 – EXPANSIONSRISIKEN (Duo B)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "DUO (TEIL B) — EXPANSIONSRISIKEN", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Was bei der Expansion schiefgehen kann", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

risks = [
    ("Qualitätserosion", "Mehr Volumen → weniger Sorgfalt. Fehlerhafte Pakete, Verzögerungen — Vertrauensverlust ist schwer reversibel."),
    ("Lieferanten-\nabhängigkeit", "Nur ein Lieferant = hohes Risiko. Ausfall wegen Insolvenz, Rohstoffmangel oder Großkunden-Priorität."),
    ("Saisonale\nSchwankungen", "Weihnachten: Hochbetrieb. Januar: Einbruch. Personal für Spitze → zu viel im Tief. Planung ist alles."),
    ("Kapitalfehler", "Zu früh zu viel investieren. Maschinen, Lager, Personal — bevor Erträge kommen. Keine Reserve = gefährlich."),
    ("Wettbewerb", "Streuselmarkt ist kein Nischenmarkt mehr. Wer nicht innoviert, verliert seinen Vorsprung an schnellere Anbieter."),
]

for i, (name, desc) in enumerate(risks):
    col = i % 3
    row = i // 3
    x = 1.0 + col * 4.0
    y = 1.7 + row * 2.4
    add_rect(slide, x, y, 3.6, 0.45, fill_color=PINK)
    add_text(slide, name, x + 0.15, y + 0.08, 3.3, 0.4, font_size=Pt(11), bold=True, color=WHITE)
    add_rect(slide, x, y + 0.45, 3.6, 1.8, fill_color=LIGHT, line_color=BORDER, line_width=Pt(1))
    add_text(slide, desc, x + 0.2, y + 0.6, 3.2, 1.5, font_size=Pt(10), color=GRAY)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Duo · Superstreusel UG · Wege in die Selbstständigkeit 2026", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY)

# ──────────────────────────────────────────────
# SLIDE 9 – FAZIT
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 0.5, 7.5, fill_color=PINK)
add_rect(slide, 0.5, 0, 12.83, 7.5, fill_color=WHITE)
add_rect(slide, 0.5, 0, 12.83, 1.4, fill_color=LIGHT)

add_text(slide, "FAZIT & AUSBLICK", 1.0, 0.2, 10, 0.4, font_size=Pt(9), bold=True, color=PINK)
add_text(slide, "Was Superstreusel uns zeigt", 1.0, 0.6, 11, 0.6, font_size=Pt(22), bold=True, color=DARK)

learnings = [
    ("Rechtsform", "Die UG ermöglicht partnerschaftliche Gründung mit Haftungsschutz — ohne das hohe Kapital der GmbH."),
    ("Finanzierung", "Innenfinanzierung funktioniert bei organischem Wachstum. Kein Kredit ist eine Philosophie, keine Schwäche."),
    ("Personal", "Zu früh warten kostet Qualität. Erster Schritt: Minijob Versand. Schrittweise skalieren."),
    ("Lieferanten", "Messen, Netzwerk, B2B-Plattformen — immer mindestens zwei Quellen aufbauen."),
    ("Marketing", "Visuelles Produkt + Community = organische Reichweite. User Generated Content ist das effizienteste Medium."),
    ("Risiken", "Qualität, Abhängigkeiten, Saisonalität — wer sie kennt, kann sie steuern."),
]

for i, (topic, text) in enumerate(learnings):
    col = i % 2
    row = i // 2
    x = 1.0 + col * 6.0
    y = 1.7 + row * 1.55
    add_rect(slide, x, y, 0.06, 0.9, fill_color=PINK)
    add_text(slide, topic, x + 0.25, y, 2.5, 0.35, font_size=Pt(11), bold=True, color=DARK)
    add_text(slide, text, x + 0.25, y + 0.4, 5.4, 0.6, font_size=Pt(10), color=GRAY)

# Final quote
add_rect(slide, 1.0, 6.3, 11.5, 0.5, fill_color=LIGHT)
add_text(slide, '"Aus einem Problem wird eine Idee - aus einer Idee wird Superstreusel."', 1.5, 6.38, 10.5, 0.35, font_size=Pt(11), italic=True, color=PINK, align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 6.9, 12.83, 0.6, fill_color=LIGHT)
add_text(slide, "Superstreusel UG (haftungsbeschränkt) · Wege in die Selbstständigkeit 2026 · superstreusel.de", 1.0, 6.95, 11, 0.4, font_size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "Superstreusel_Praesentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
